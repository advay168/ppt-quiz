from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.enum.dml import MSO_THEME_COLOR
from random import shuffle
from lxml import etree
from string import ascii_uppercase

def hideSlide(slide):
    sld = slide._element
    sld.set('show', '0')


def hyperlink(run_object, source_slide, destination_slide):
    rId = source_slide.part.relate_to(destination_slide.part, RT.SLIDE)
    rPr = run_object._r.get_or_add_rPr()
    hlinkClick = rPr.add_hlinkClick(rId)
    hlinkClick.set('action', 'ppaction://hlinksldjump')


def setTitle(slide, title):
    slide.shapes.title.text = title


def add_textbox(top, left, width, height, slide):
    return slide.shapes.add_textbox(left, top, width, height)

def removeBullets(para):
    para._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"))


def makePPTFromQs(Questions_and_answers, name):
    prs = Presentation()

    title_slide_layout = prs.slide_layouts[0]
    title_content_slide_layout = prs.slide_layouts[1]
    title_only_slide_layout = prs.slide_layouts[5]

    title_slide = prs.slides.add_slide(title_slide_layout)
    setTitle(title_slide, "Quiz")

    quiz_slides = []

    correct_slides = []
    incorrect_slides = []
    c_next_runs = []
    ic_next_runs = []

    for index,Q in enumerate(Questions_and_answers):
        current_slide = prs.slides.add_slide(title_content_slide_layout)
        setTitle(current_slide,f"Q{index+1} {Q['Question']}")
        quiz_slides.append(current_slide)

        correct_slide = prs.slides.add_slide(title_only_slide_layout)
        setTitle(correct_slide, "Correct!")
        correct_slides.append(correct_slide)
        tb = add_textbox(prs.slide_height-Inches(1.15), prs.slide_width -
                         Inches(2), Inches(1), Inches(1), correct_slide)
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = "Next"
        r.font.size=Pt(25)
        r.font.color.theme_color=MSO_THEME_COLOR.ACCENT_4
        c_next_runs.append(r)

        incorrect_slide = prs.slides.add_slide(title_content_slide_layout)
        setTitle(incorrect_slide, "Incorrect!")
        incorrect_slides.append(incorrect_slide)
        r=incorrect_slide.shapes[1].text_frame.paragraphs[0].add_run()
        r.text=f"The correct answer is {Q['Answer']}"
        r.font.color.theme_color=MSO_THEME_COLOR.ACCENT_4
        r.font.size=Pt(25)
        removeBullets(incorrect_slide.shapes[1].text_frame.paragraphs[0])
        tb = add_textbox(prs.slide_height-Inches(1.15), prs.slide_width -
                         Inches(2), Inches(1), Inches(1), incorrect_slide)
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = "Next"
        r.font.size=Pt(25)
        r.font.color.theme_color=MSO_THEME_COLOR.ACCENT_4
        ic_next_runs.append(r)

        options = Q["Options"]+[Q["Answer"]]
        shuffle(options)
        para = current_slide.shapes[1].text_frame.paragraphs[0]
        removeBullets(para)
        for i, option in enumerate(options):
            new_line = "" if i == len(options)-1 else "\n"
            run = para.add_run()
            run.text = ascii_uppercase[i%26] + ") " +option+new_line
            run.font.color.theme_color=MSO_THEME_COLOR.ACCENT_5
            run.font.size=Pt(30)
            if option == Q["Answer"]:
                hyperlink(run, current_slide, correct_slide)
            else:
                hyperlink(run, current_slide, incorrect_slide)
    
    start_run=title_slide.shapes[1].text_frame.paragraphs[0].add_run()
    start_run.text="Start"
    hyperlink(start_run,title_slide,quiz_slides.pop(0))
    
    thank_slide=prs.slides.add_slide(title_only_slide_layout)
    thank_slide.shapes.title.text="Thank you for attempting"
    quiz_slides.append(thank_slide)
    for r1, r2, c_slide, ic_slide, q_slide in zip(c_next_runs, ic_next_runs, correct_slides, incorrect_slides, quiz_slides):
        hyperlink(r1,c_slide,q_slide)
        hyperlink(r2,ic_slide,q_slide)

    prs.save(name)

if __name__ == "__main__":
    makePPTFromQs([
        {"Question":"Q1","Answer":"A31","Options":["O1","O2"]},{"Question":"Q2","Answer":"A2","Options":["O1","O2"]},{"Question":"Q3","Answer":"A3","Options":["O1","O2"]}
    ],"Test1.pptx")